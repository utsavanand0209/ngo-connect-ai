#!/usr/bin/env python3
"""
NGO-Connect: Stakeholder Presentation Generator
Generates a polished PowerPoint presentation from project facts.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x0B, 0x1D, 0x51)
TEAL      = RGBColor(0x00, 0x96, 0x88)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
GREY_TEXT  = RGBColor(0x55, 0x55, 0x55)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_BLUE   = RGBColor(0x1E, 0x88, 0xE5)
ACCENT_GREEN  = RGBColor(0x43, 0xA0, 0x47)
ACCENT_PURPLE = RGBColor(0x7B, 0x1F, 0xA2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg_rect(slide, color):
    """Full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        SLIDE_W, SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_accent_bar(slide, y=Inches(0), height=Inches(0.06), color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), y,
        SLIDE_W, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=DARK_TEXT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_stat_card(slide, left, top, width, height, number, label, color):
    """Rounded-look stat card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREY_TEXT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER


def add_section_header(slide, title, subtitle=""):
    add_bg_rect(slide, NAVY)
    add_accent_bar(slide, y=Inches(3.5), height=Inches(0.04), color=TEAL)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
                 title, 44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.8),
                     subtitle, 20, color=RGBColor(0xB0, 0xBE, 0xC5),
                     alignment=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────
# BUILD PRESENTATION
# ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, NAVY)
add_accent_bar(sl, y=Inches(4.8), height=Inches(0.05), color=TEAL)

add_text_box(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
             "NGO-Connect", 56, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.3), Inches(11), Inches(1),
             "Design & Architecture Overview", 28, color=TEAL,
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(5.2), Inches(11), Inches(0.8),
             "Connecting NGOs • Donors • Volunteers • Administrators",
             18, color=RGBColor(0xB0, 0xBE, 0xC5), alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
             "February 2026  |  Stakeholder Briefing",
             14, color=RGBColor(0x78, 0x90, 0x9C), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Agenda", 36, bold=True, color=NAVY)

agenda_items = [
    "1.  Platform Overview & Value Proposition",
    "2.  System Architecture & Technology Stack",
    "3.  User Roles & Access Control",
    "4.  Core Feature Modules",
    "5.  Database Design",
    "6.  REST API Surface (89 Endpoints)",
    "7.  Key Workflows (Donations, Volunteering, Moderation)",
    "8.  AI & Recommendation Engine",
    "9.  Frontend Application Structure",
    "10. Deployment & Configuration",
]
add_bullet_list(sl, Inches(1.2), Inches(1.6), Inches(10), Inches(5.5),
                agenda_items, font_size=20, color=DARK_TEXT, spacing=Pt(10))

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 – EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Executive Summary", 36, bold=True, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
             "NGO-Connect is a full-stack web platform bridging NGOs, donors, volunteers, and "
             "administrators. It enables transparent donations with payment processing, volunteer "
             "management with certificate issuance, AI-powered recommendations, and comprehensive "
             "platform administration.",
             17, color=DARK_TEXT)

# Stat cards
add_stat_card(sl, Inches(0.8), Inches(3.2), Inches(2.5), Inches(1.8),
              "89", "REST API\nEndpoints", ACCENT_BLUE)
add_stat_card(sl, Inches(3.7), Inches(3.2), Inches(2.5), Inches(1.8),
              "13", "Route\nModules", TEAL)
add_stat_card(sl, Inches(6.6), Inches(3.2), Inches(2.5), Inches(1.8),
              "17", "Database\nTables", ACCENT_ORANGE)
add_stat_card(sl, Inches(9.5), Inches(3.2), Inches(2.5), Inches(1.8),
              "28+", "Frontend\nPages", ACCENT_PURPLE)

add_bullet_list(sl, Inches(0.8), Inches(5.4), Inches(11.5), Inches(2),
                [
                    "✓  Real payment integration (Razorpay + Mock mode)",
                    "✓  AI chatbot powered by Google Gemini with RAG",
                    "✓  Role-based access: User, NGO, Admin",
                    "✓  Certificate generation for donations & volunteering",
                ],
                font_size=16, color=GREY_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 – SECTION: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "System Architecture", "Three-tier design with PostgreSQL backbone")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 – HIGH-LEVEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "High-Level Architecture", 36, bold=True, color=NAVY)

# Tier boxes
tiers = [
    ("Presentation Tier", "React 18 SPA  •  Tailwind CSS  •  Recharts  •  Leaflet Maps\nAxios + JWT Auth  •  React Router v6", ACCENT_BLUE),
    ("Application Tier", "Node.js + Express  •  13 Route Modules  •  JWT Middleware\nPayment Gateway  •  Gemini AI  •  Certificate Engine", TEAL),
    ("Data Tier", "PostgreSQL 14+  •  13 Primary Tables + 4 Junction Tables\nBIGSERIAL PK + UUID External ID  •  JSONB source_doc  •  Relational FKs", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(tiers):
    y = Inches(1.6) + Inches(i * 1.85)
    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), y, Inches(10), Inches(1.55),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2.5)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = GREY_TEXT
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

# Arrow indicators between tiers
for i in range(2):
    y = Inches(3.15) + Inches(i * 1.85)
    add_text_box(sl, Inches(5.8), y, Inches(1.5), Inches(0.5),
                 "▼", 24, bold=True, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# External services note
add_text_box(sl, Inches(1), Inches(7), Inches(11), Inches(0.4),
             "External Services:  Razorpay Payment API  •  Google Gemini API (gemini-2.5-flash)",
             13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 – TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Technology Stack", 36, bold=True, color=NAVY)

# Backend column
add_text_box(sl, Inches(1), Inches(1.5), Inches(5), Inches(0.5),
             "Backend", 24, bold=True, color=TEAL)
backend_tech = [
    "Node.js + Express ^4.22",
    "PostgreSQL 14+ (pg ^8.16)",
    "JWT Auth (jsonwebtoken ^9.0)",
    "bcryptjs ^2.4 (password hashing)",
    "Multer ^1.4 (file uploads)",
    "Google Generative AI ^0.24",
    "Razorpay / Mock payments",
    "dotenv + CORS",
]
add_bullet_list(sl, Inches(1.2), Inches(2.1), Inches(5), Inches(4.5),
                backend_tech, font_size=16, color=DARK_TEXT, spacing=Pt(8))

# Frontend column
add_text_box(sl, Inches(7), Inches(1.5), Inches(5), Inches(0.5),
             "Frontend", 24, bold=True, color=ACCENT_BLUE)
frontend_tech = [
    "React ^18.2 + React Router v6",
    "Axios ^1.4 (HTTP client)",
    "Tailwind CSS ^3.4",
    "Recharts ^2.8 (data viz)",
    "Leaflet ^1.9 + react-leaflet ^4.2",
    "leaflet-routing-machine ^3.2",
    "Heroicons ^1.0 (icons)",
    "react-scripts ^5.0 (CRA toolchain)",
]
add_bullet_list(sl, Inches(7.2), Inches(2.1), Inches(5), Inches(4.5),
                frontend_tech, font_size=16, color=DARK_TEXT, spacing=Pt(8))

# Divider
shape = sl.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(0.03), Inches(4.5),
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
shape.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 – USER ROLES & ACCESS CONTROL
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "User Roles & Access Control", 36, bold=True, color=NAVY)

roles = [
    ("User (Donor / Volunteer)", ACCENT_BLUE,
     ["Discover & donate to NGO campaigns", "Apply for volunteer opportunities",
      "Submit help requests to NGOs", "View certificates & receipts",
      "Message NGOs, get AI recommendations"]),
    ("NGO", TEAL,
     ["Manage organization profile & verification docs", "Create campaigns & volunteer opportunities",
      "Review & approve certificate requests", "Process help requests from users",
      "Message users & manage volunteers"]),
    ("Admin", ACCENT_ORANGE,
     ["Verify / reject NGO registrations", "Moderate flagged content & flag requests",
      "Manage categories & broadcast notifications",
      "Platform analytics & dashboard (JSON + SSR HTML)",
      "User management (delete users, enable/disable NGOs)"]),
]

for i, (title, color, items) in enumerate(roles):
    left = Inches(0.5) + Inches(i * 4.2)
    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6),
        Inches(3.9), Inches(5.2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    add_text_box(sl, left + Inches(0.3), Inches(1.8), Inches(3.3), Inches(0.6),
                 title, 18, bold=True, color=color)
    add_bullet_list(sl, left + Inches(0.3), Inches(2.5), Inches(3.3), Inches(4),
                    items, font_size=13, color=DARK_TEXT, spacing=Pt(6))

add_text_box(sl, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
             "Auth: Stateless JWT (7-day expiry)  •  bcrypt password hashing  •  Role-checked middleware on every route",
             13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 – SECTION: CORE FEATURES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "Core Feature Modules",
                   "Donations • Volunteering • Messaging • Moderation • AI")

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 – FEATURE MODULES OVERVIEW
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Feature Modules Overview", 36, bold=True, color=NAVY)

features = [
    ("Donations", "Campaign-based donations with Razorpay/mock payments,\nreceipt generation, and NGO-approved certificates", ACCENT_BLUE),
    ("Volunteering", "Standalone opportunities + campaign volunteering,\napplication lifecycle, activity completion & certificates", TEAL),
    ("Messaging", "User-to-NGO and NGO-to-user messaging with\nconversation threads, unread counts, and broadcast", ACCENT_ORANGE),
    ("Help Requests", "Users submit support requests to NGOs with\nstatus workflow: Pending → Approved → Completed", ACCENT_GREEN),
    ("Moderation", "User flag requests for NGOs/campaigns,\nadmin review & resolution, enable/disable NGOs", ACCENT_PURPLE),
    ("AI Intelligence", "Personalized recommendations, LLM chatbot with RAG,\ncampaign classification, fraud scoring, volunteer matching", RGBColor(0xE5, 0x39, 0x35)),
]

for i, (title, desc, color) in enumerate(features):
    row = i // 3
    col = i % 3
    left = Inches(0.5) + Inches(col * 4.2)
    top = Inches(1.5) + Inches(row * 2.8)

    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
        Inches(3.9), Inches(2.3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    add_text_box(sl, left + Inches(0.25), top + Inches(0.25),
                 Inches(3.4), Inches(0.5),
                 title, 20, bold=True, color=color)
    add_text_box(sl, left + Inches(0.25), top + Inches(0.85),
                 Inches(3.4), Inches(1.2),
                 desc, 13, color=GREY_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 – SECTION: DATABASE
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "Database Design",
                   "Hybrid Document-Relational Pattern on PostgreSQL")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 – DATABASE TABLES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Database Schema (17 Tables)", 36, bold=True, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(1.3), Inches(11), Inches(0.7),
             "Design: BIGSERIAL PK + UUID external_id + Typed relational columns + JSONB source_doc + Foreign Keys + Indexes",
             15, color=GREY_TEXT)

# Primary tables
add_text_box(sl, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.5),
             "Primary Tables (13)", 20, bold=True, color=TEAL)

primary_tables = [
    "users_rel — Platform users (donor, volunteer, admin)",
    "ngos_rel — NGO organizations",
    "categories_rel — NGO/campaign categories",
    "campaigns_rel — Fundraising campaigns → ngos_rel",
    "donations_rel — Financial contributions → users, campaigns, ngos",
    "volunteer_opportunities_rel — Volunteer positions → ngos",
    "volunteer_applications_rel — Applications → users, opportunities",
    "certificates_rel — Donation & volunteer certificates",
    "messages_rel — User ↔ NGO messaging",
    "notifications_rel — Platform notifications",
    "help_requests_rel — Support requests → users, ngos",
    "flag_requests_rel — Content moderation flags",
    "ai_logs_rel — AI operation audit trail",
]
add_bullet_list(sl, Inches(1), Inches(2.7), Inches(6.5), Inches(4.5),
                primary_tables, font_size=12, color=DARK_TEXT, spacing=Pt(3))

# Junction tables
add_text_box(sl, Inches(8), Inches(2.1), Inches(4.5), Inches(0.5),
             "Junction Tables (4)", 20, bold=True, color=ACCENT_ORANGE)

junction_tables = [
    "ngo_categories_rel\n  NGO ↔ Category names",
    "campaign_volunteers_rel\n  Campaign ↔ User volunteers",
    "campaign_volunteer_registrations_rel\n  Campaign sign-ups with details",
    "opportunity_applicants_rel\n  Opportunity ↔ User applicants",
]
add_bullet_list(sl, Inches(8.2), Inches(2.7), Inches(4.3), Inches(3),
                junction_tables, font_size=13, color=DARK_TEXT, spacing=Pt(8))

# Key relationships
add_text_box(sl, Inches(8), Inches(4.8), Inches(4.5), Inches(0.5),
             "Key Relationships", 20, bold=True, color=ACCENT_BLUE)
rels = [
    "campaigns → ngos (ngo_id FK)",
    "donations → users, campaigns, ngos",
    "vol. applications → users, opportunities",
    "certificates → donations | vol. applications",
    "messages: from_user ↔ to_ngo",
]
add_bullet_list(sl, Inches(8.2), Inches(5.4), Inches(4.3), Inches(2),
                rels, font_size=13, color=DARK_TEXT, spacing=Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 – SECTION: API
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "REST API Surface",
                   "89 Endpoints across 13 Route Modules")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 – API ENDPOINT SUMMARY
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "API Endpoint Distribution", 36, bold=True, color=NAVY)

api_modules = [
    ("/api/admin", "18", ACCENT_ORANGE),
    ("/api/volunteering", "12", TEAL),
    ("/api/campaigns", "11", ACCENT_BLUE),
    ("/api/donations", "8", ACCENT_GREEN),
    ("/api/ngos", "7", ACCENT_PURPLE),
    ("/api/messages", "7", RGBColor(0xE5, 0x39, 0x35)),
    ("/api/ai", "6", RGBColor(0xFF, 0xB3, 0x00)),
]
api_modules_2 = [
    ("/api/categories", "5", TEAL),
    ("/api/auth", "4", ACCENT_BLUE),
    ("/api/requests", "4", ACCENT_ORANGE),
    ("/api/users", "3", ACCENT_GREEN),
    ("/api/certificates", "3", ACCENT_PURPLE),
    ("/api/notifications", "1", GREY_TEXT),
]

for i, (path, count, color) in enumerate(api_modules):
    y = Inches(1.5) + Inches(i * 0.7)
    # bar
    bar_w = int(count) / 18 * 6  # scale relative to max (18)
    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.8), y, Inches(bar_w), Inches(0.45),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_text_box(sl, Inches(0.5), y, Inches(2.2), Inches(0.45),
                 path, 14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    add_text_box(sl, Inches(2.8) + Inches(bar_w) + Inches(0.15), y,
                 Inches(0.7), Inches(0.45),
                 count, 16, bold=True, color=color)

for i, (path, count, color) in enumerate(api_modules_2):
    y = Inches(1.5) + Inches(i * 0.7)
    bar_w = int(count) / 18 * 6
    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.3), y, Inches(max(bar_w, 0.3)), Inches(0.45),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False

    add_text_box(sl, Inches(7), y, Inches(2.2), Inches(0.45),
                 path, 14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    add_text_box(sl, Inches(9.3) + Inches(max(bar_w, 0.3)) + Inches(0.15), y,
                 Inches(0.7), Inches(0.45),
                 count, 16, bold=True, color=color)

add_text_box(sl, Inches(0.5), Inches(6.6), Inches(12), Inches(0.6),
             "Auth: Public | Any Authenticated | User | NGO | Admin  •  Content-Type: application/json  •  File uploads: multipart/form-data",
             13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 – KEY API HIGHLIGHTS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Key API Highlights", 36, bold=True, color=NAVY)

col1_items = [
    "Authentication (4 endpoints)",
    "  POST /auth/register — User or NGO signup",
    "  POST /auth/login — JWT token (7-day expiry)",
    "  GET /auth/me — Current user profile",
    "",
    "Donations (8 endpoints)",
    "  POST /donations/campaign/:id/initiate",
    "  POST /donations/:id/confirm",
    "  POST /donations/:id/certificate/decision",
    "  GET /donations/:id/receipt",
]
add_bullet_list(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                col1_items, font_size=14, color=DARK_TEXT, spacing=Pt(4))

col2_items = [
    "AI & Intelligence (6 endpoints)",
    "  GET /ai/recommendations — Personalized",
    "  POST /ai/chat — LLM chatbot + RAG",
    "  POST /ai/classify-campaign — Auto-classify",
    "  POST /ai/fraud-score — Risk scoring",
    "",
    "Admin Dashboard (18 endpoints)",
    "  GET /admin/dashboard — KPI snapshot JSON",
    "  GET /admin/dashboard/ssr — Full HTML render",
    "  GET /admin/analytics — Charts data",
]
add_bullet_list(sl, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                col2_items, font_size=14, color=DARK_TEXT, spacing=Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 – SECTION: WORKFLOWS
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "Key Workflows",
                   "Donation • Volunteer • Help Request • Moderation")

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 – DONATION WORKFLOW
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Donation Workflow", 36, bold=True, color=NAVY)

steps = [
    ("1", "User selects campaign\n& payment method", ACCENT_BLUE),
    ("2", "POST /initiate\ncreates payment order", TEAL),
    ("3", "Payment gateway\nprocesses payment", ACCENT_ORANGE),
    ("4", "POST /confirm\nverifies signature", ACCENT_GREEN),
    ("5", "Receipt generated\ncampaign updated", ACCENT_PURPLE),
    ("6", "NGO reviews &\napproves certificate", RGBColor(0xE5, 0x39, 0x35)),
]

for i, (num, desc, color) in enumerate(steps):
    left = Inches(0.3) + Inches(i * 2.1)
    # Circle number
    circ = sl.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.65), Inches(1.8),
        Inches(0.7), Inches(0.7),
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()

    tf = circ.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(sl, left + Inches(0.1), Inches(2.7),
                 Inches(1.8), Inches(1.2),
                 desc, 13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        add_text_box(sl, left + Inches(1.7), Inches(1.9),
                     Inches(0.5), Inches(0.5),
                     "→", 24, bold=True, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
             "Payment Methods: UPI • Credit/Debit Card • Net Banking  |  Gateways: Razorpay (prod) • Mock (dev)",
             15, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# Volunteer workflow
add_text_box(sl, Inches(0.8), Inches(5.2), Inches(11), Inches(0.6),
             "Volunteer Workflow", 28, bold=True, color=NAVY)

vol_steps = [
    "Browse opportunities / campaign volunteer roles →",
    "Apply with contact details & motivation →",
    "Complete activity & log hours →",
    "Request certificate →",
    "NGO approves → Certificate issued"
]
add_bullet_list(sl, Inches(1), Inches(5.9), Inches(11), Inches(1.5),
                vol_steps, font_size=14, color=DARK_TEXT, spacing=Pt(4))

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 – SECTION: AI
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "AI & Recommendation Engine",
                   "Gemini LLM • Rule-Based Scoring • RAG Pipeline")

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 – AI FEATURES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "AI & Intelligence Features", 36, bold=True, color=NAVY)

ai_features = [
    ("Personalized\nRecommendations", "Rule-based scoring using user preferences\n(location, interests, causes, skills)\nmatched against NGO sectors & campaigns.\nTop 10 each with scores + reasons.", ACCENT_BLUE),
    ("LLM Chatbot\nwith RAG", "Google Gemini (gemini-2.5-flash)\n+ DB retrieval of relevant NGOs/campaigns\n+ 13-article Knowledge Base\n+ Role-aware prompts + fallback mode", TEAL),
    ("Campaign\nClassification", "Keyword-based auto-categorization:\nEducation, Health, Food,\nDisaster Relief, Environment, Other", ACCENT_ORANGE),
    ("Fraud Scoring", "Heuristic analysis: verification docs,\naccount age, suspicious keywords,\nunrealistic goals. Flags if score ≥ 50.", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(ai_features):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + Inches(col * 6.3)
    top = Inches(1.5) + Inches(row * 2.7)

    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
        Inches(5.9), Inches(2.3),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    add_text_box(sl, left + Inches(0.3), top + Inches(0.2),
                 Inches(5.3), Inches(0.7),
                 title, 18, bold=True, color=color)
    add_text_box(sl, left + Inches(0.3), top + Inches(0.9),
                 Inches(5.3), Inches(1.2),
                 desc, 13, color=DARK_TEXT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 – SECTION: FRONTEND
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_section_header(sl, "Frontend Application",
                   "React 18 SPA with 28+ Pages & Role-Based Routing")

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 – FRONTEND PAGES
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Frontend Page Architecture", 36, bold=True, color=NAVY)

page_groups = [
    ("Public (8)", [
        "Home — Landing with hero CTA, helplines, categories",
        "Login / Register — Dual-mode (user + NGO)",
        "NGO List & Profile — Search, filter, tabs, charts",
        "Campaign List & Details — Donation + volunteer flows",
        "Chatbot — AI support with suggested questions",
    ], ACCENT_BLUE),
    ("User (7)", [
        "Dashboard — Help requests, donations, volunteers, certs",
        "Donate — Campaign selection, payment, receipts",
        "Volunteer Opportunities — Apply, withdraw, complete",
        "Recommendations — AI-powered with match scores",
        "Messages / Profile / Map — Communication & navigation",
    ], TEAL),
    ("Admin (7)", [
        "Dashboard — KPI cards, charts, auto-refresh, SSR viewer",
        "Verifications — NGO approval/rejection queue",
        "Flagged Content — Moderation & flag request review",
        "Analytics — Line/bar charts, platform statistics",
        "Users / Notifications / Categories / Requests",
    ], ACCENT_ORANGE),
]

for i, (title, items, color) in enumerate(page_groups):
    left = Inches(0.4) + Inches(i * 4.2)

    add_text_box(sl, left, Inches(1.5), Inches(3.9), Inches(0.5),
                 title, 20, bold=True, color=color)

    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.1),
        Inches(3.9), Inches(4.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False

    add_bullet_list(sl, left + Inches(0.2), Inches(2.3),
                    Inches(3.5), Inches(4),
                    items, font_size=12, color=DARK_TEXT, spacing=Pt(8))

add_text_box(sl, Inches(0.5), Inches(6.9), Inches(12), Inches(0.5),
             "Components: Navbar • ProtectedRoute • UserRoute • AdminRoute • ConfirmModal • PreferencesModal • RecommendedNgos",
             13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 21 – DEPLOYMENT & CONFIG
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, WHITE)
add_accent_bar(sl, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Deployment & Configuration", 36, bold=True, color=NAVY)

add_text_box(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "Backend Environment", 22, bold=True, color=TEAL)
env_vars = [
    "PORT = 5001",
    "POSTGRES_URL = postgresql://...",
    "JWT_SECRET = <strong-secret>",
    "GEMINI_API_KEY = <optional>",
    "PAYMENT_GATEWAY_PROVIDER = mock | razorpay",
    "RAZORPAY_KEY_ID = <optional>",
    "RAZORPAY_KEY_SECRET = <optional>",
]
add_bullet_list(sl, Inches(1), Inches(2.1), Inches(5.5), Inches(3),
                env_vars, font_size=14, color=DARK_TEXT, spacing=Pt(5))

add_text_box(sl, Inches(7), Inches(1.5), Inches(5), Inches(0.5),
             "Setup Commands", 22, bold=True, color=ACCENT_BLUE)
commands = [
    "npm install",
    "npm run db:relational-schema",
    "npm run seed",
    "npm run dev (backend)",
    "npm start (frontend, port 3000)",
    "npm run smoke (API test)",
]
add_bullet_list(sl, Inches(7.2), Inches(2.1), Inches(5), Inches(3),
                commands, font_size=14, color=DARK_TEXT, spacing=Pt(5))

add_text_box(sl, Inches(0.8), Inches(5), Inches(11), Inches(0.5),
             "Seed Credentials", 22, bold=True, color=ACCENT_ORANGE)
creds = [
    "Admin:  admin@ngoconnect.org  /  password123",
    "User:   rahul@example.com  /  password123",
    "NGO:    akshayapatra@ngo.org  /  password123",
]
add_bullet_list(sl, Inches(1), Inches(5.6), Inches(11), Inches(1.5),
                creds, font_size=15, color=DARK_TEXT, spacing=Pt(6))

# ═══════════════════════════════════════════════════════════════
# SLIDE 22 – THANK YOU
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg_rect(sl, NAVY)
add_accent_bar(sl, y=Inches(4.6), height=Inches(0.05), color=TEAL)

add_text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
             "Thank You", 56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "Questions & Discussion", 28, color=TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(1), Inches(5.0), Inches(11), Inches(1),
             "89 Endpoints  •  17 Tables  •  28+ Pages  •  AI-Powered\n"
             "Full details: DESIGN_AND_ARCHITECTURE.md",
             16, color=RGBColor(0xB0, 0xBE, 0xC5), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "/Users/ppushkar/Development/experiment/ngo-connect-ai/NGO_Connect_Architecture_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
